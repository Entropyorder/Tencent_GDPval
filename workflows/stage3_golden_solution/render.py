#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
golden-solution-pro renderer.

输入: content.json (内容层 claude -p 产出) + 风格包
输出: .docx / .xlsx / .pptx / .md

铁律: 本脚本只渲染,不生成内容。所有文字来自 content.json。
黑白强制: _is_grey() 守卫拦截一切蓝色/彩色,统一映射为黑/灰。
"""
import sys, os, json, argparse, re

# ---------- 依赖自检 ----------
def _ensure_deps():
    missing = []
    try: import docx
    except ImportError: missing.append("python-docx")
    try: import openpyxl
    except ImportError: missing.append("openpyxl")
    if missing:
        print(f"[render] 缺失依赖 {missing}, 自动安装到用户目录...", file=sys.stderr)
        import subprocess
        subprocess.check_call([sys.executable, "-m", "pip", "install", "--user", *missing])
        # 重新 import
        import importlib
        if "python-docx" in missing:
            global docx; import docx; importlib.reload(docx)
        if "openpyxl" in missing:
            import openpyxl; importlib.reload(openpyxl)
_ensure_deps()
import docx
from docx.shared import Pt, Cm, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

# ======================================================================
# 风格包
# ======================================================================
# 全部黑白。颜色只允许: 黑 #000000 / 深灰 #404040 / 中灰 #808080 / 浅灰 #D9D9D9 / 极浅灰 #F2F2F2 / 白 #FFFFFF
GREY_BLACK = RGBColor(0x00, 0x00, 0x00)
GREY_DARK  = RGBColor(0x40, 0x40, 0x40)
GREY_MID   = RGBColor(0x80, 0x80, 0x80)
GREY_LIGHT = RGBColor(0xD9, 0xD9, 0xD9)
GREY_PALE  = RGBColor(0xF2, 0xF2, 0xF2)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

# Excel 端 openpyxl 用 hex 串
XL_BLACK = "000000"; XL_DARK="404040"; XL_MID="808080"; XL_LIGHT="D9D9D9"; XL_PALE="F2F2F2"; XL_WHITE="FFFFFF"

# ---------- ref-semi 参考风格真实色值（彩色，复刻某半导体 2026Q1 真实底稿）----------
# 仅用于 word-ref / excel-ref 包；A/B/C 黑白包不受影响。
REF_HEADER_BLUE   = "1F4E78"   # 深蓝：表头底纹
REF_INPUT_YELLOW  = "FFF2CC"   # 黄：输入（直接取自附件原值）
REF_COMPUTE_GREEN = "E2EFDA"   # 绿：公式/交叉计算
REF_CAUTION_ORANGE= "FCE4D6"   # 橙：提示/注意/风险得分
REF_RISK_ORANGE   = "F8CBAD"   # 浅橙：风险得分高亮（=概率×影响）
REF_TITLE_BLUE_RGB= RGBColor(0x1F, 0x4E, 0x78)   # 标题深蓝字色
REF_WHITE_RGB     = RGBColor(0xFF, 0xFF, 0xFF)
REF_BLACK_RGB     = RGBColor(0x00, 0x00, 0x00)
# 数据性质 → 底纹 hex（Excel/Word 共用）
REF_KIND_FILL = {
    "input":   REF_INPUT_YELLOW,
    "compute": REF_COMPUTE_GREEN,
    "caution": REF_CAUTION_ORANGE,
    "risk":    REF_RISK_ORANGE,
}

def _is_grey(rgb):
    """守卫: 颜色非灰阶一律降级为黑或灰。蓝/绿/红一律拒绝。"""
    if rgb is None: return GREY_BLACK
    if isinstance(rgb, str):
        if not rgb.startswith("#") or len(rgb) != 7: return GREY_BLACK
        r = int(rgb[1:3],16); g = int(rgb[3:5],16); b = int(rgb[5:7],16)
    elif isinstance(rgb, (list,tuple)) and len(rgb)>=3:
        r,g,b = rgb[0],rgb[1],rgb[2]
    else:
        return GREY_BLACK
    mx,mn = max(r,g,b), min(r,g,b)
    # 非灰阶判定: 三通道差值过大 (彩色)
    if mx - mn > 20:
        return GREY_BLACK   # 任何彩色 → 黑
    v = (r+g+b)//3
    # 映射到标准灰阶档
    if v >= 235: return WHITE
    if v >= 190: return GREY_PALE
    if v >= 150: return GREY_LIGHT
    if v >= 90:  return GREY_MID
    if v >= 50:  return GREY_DARK
    return GREY_BLACK

# ---------- Word 风格包 ----------
WORD_PACKS = {
    "word-A": {  # 庄重公文
        "body_font": "宋体", "body_size": 12, "line_spacing": 1.5,
        "indent_first_line": True,  # 首行缩进2字
        "heading_font": "黑体", "heading_color": GREY_BLACK,
        "heading_scheme": "cn",   # 一、二、三
        "caption_style": "表N", "caption_sep": "　",  # 全角空格
        "footnote_label": "注：",
        "para_space_after": 0,
    },
    "word-B": {  # 现代咨询
        "body_font": "等线", "body_size": 10.5, "line_spacing": 1.15,
        "indent_first_line": False,
        "heading_font": "等线", "heading_color": GREY_DARK,
        "heading_scheme": "1.1",
        "caption_style": "表 N", "caption_sep": " ",
        "footnote_label": "数据来源：",
        "para_space_after": 6,
    },
    "word-C": {  # 学术尽调
        "body_font": "宋体", "body_size": 11, "line_spacing": 1.4,
        "indent_first_line": True,
        "heading_font": "Times New Roman", "heading_color": GREY_BLACK,
        "heading_scheme": "1.1",
        "caption_style": "Table N", "caption_sep": " ",
        "footnote_label": "注：",
        "para_space_after": 0,
        "threeline_table": True,
    },
    # ---- 参考风格：复刻某半导体 2026Q1 合规与财务分析报告（彩色）----
    "word-ref": {
        "allow_color": True,            # 跳过 _is_grey 守卫，使用真实色值
        "body_font": "宋体", "body_size": 10.5, "fine_size": 9,
        "line_spacing": 1.0, "indent_first_line": False,
        "heading_font": "宋体", "heading_color": REF_BLACK_RGB,
        "heading_bold": False,          # 标题不加粗
        "heading_scheme": "cn+1.1",     # 一级"一、"，二级"1.1"
        "h1_size": 15, "h2_size": 12.5,
        "title_size": 20, "subtitle_size": 16,
        "caption_style": "表N", "caption_sep": "　",
        "footnote_label": "数据来源：",
        "para_space_after": 0,
        "table_full_grid": True,        # 全网格 0.5pt 黑，非三线表
        "header_fill": REF_HEADER_BLUE, # 表头深蓝
        "header_font_color": REF_WHITE_RGB,
        "header_size": 9.5,
    },
}

# ---------- Excel 风格包 ----------
EXCEL_PACKS = {
    "excel-A": {  # 经典报表
        "title_merge": True, "header_fill": XL_LIGHT, "header_bold": True,
        "border": "thin", "freeze_first_row": True, "thousands_sep": True,
        "caption_above": True, "caption_label": "表N",
        "total_row": "bold_top_double",  # 合计行: 上加粗+双线
    },
    "excel-B": {  # 紧凑仪表
        "title_merge": True, "header_fill": XL_DARK, "header_bold": True,
        "header_font_color": XL_WHITE,
        "border": "none", "freeze_first_row": True, "thousands_sep": True,
        "cond_color": True,  # 正数深灰、负数黑(非红/绿彩),用粗体区分
        "caption_above": True, "caption_label": "表 N",
    },
    "excel-C": {  # 审计底稿
        "title_merge": False, "header_fill": None, "header_bold": True,
        "border": "threeline", "freeze_first_row": True, "thousands_sep": True,
        "currency_right": True, "caption_above": True, "caption_label": "表 N",
        "total_row": "bold_top_double",
        "cell_comment_source": True,
    },
    # ---- 参考风格：复刻某半导体 2026Q1 核查底稿/风险台账（彩色）----
    "excel-ref": {
        "allow_color": True,
        "title_merge": True, "title_size": 13, "title_color": REF_HEADER_BLUE,
        "title_row_height": 22,
        "desc_row": True,                 # 行2 说明行
        "header_fill": REF_HEADER_BLUE, "header_font_color": XL_WHITE,
        "header_bold": True, "header_size": 11,
        "border": "thin", "freeze_first_row": False,  # 参考文件未冻结
        "thousands_sep": False,          # 由 col_fmt 显式控制
        "caption_above": True, "caption_label": "表N",
        "data_font": "Calibri", "data_size": 11,
        "row_height": 46,                 # 数据行加高
        "kind_fill": REF_KIND_FILL,      # 数据性质色码
        "legend": True,                  # 说明表底部图例
    },
}

# ======================================================================
# 工具: 段落/表格黑白加固
# ======================================================================
def _force_run_black(run, color=None):
    run.font.color.rgb = _is_grey(color) if color else GREY_BLACK

def _force_cell_shade(cell, fill_hex):
    """强制单元格底纹为灰阶。"""
    v = _is_grey(f"#{fill_hex}" if fill_hex else None)
    # 转 hex
    hexs = "%02X%02X%02X" % (v[0], v[1], v[2]) if hasattr(v,"__getitem__") and not isinstance(v,RGBColor) else v
    if isinstance(v, RGBColor):
        hexs = str(v)
    tcPr = cell._tc.get_or_add_tcPr()
    # 清除已有 shd
    for el in tcPr.findall(qn('w:shd')):
        tcPr.remove(el)
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), hexs)
    tcPr.append(shd)

def _force_cell_shade_raw(cell, fill_hex):
    """直接写入真实色值底纹，不经 _is_grey 守卫（用于 ref-semi 参考风格）。"""
    tcPr = cell._tc.get_or_add_tcPr()
    for el in tcPr.findall(qn('w:shd')):
        tcPr.remove(el)
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), str(fill_hex).lstrip('#').upper())
    tcPr.append(shd)

def _set_cell_border(cell, **kwargs):
    """边框黑白。kwargs: top/left/bottom/right -> {sz, val}"""
    tcPr = cell._tc.get_or_add_tcPr()
    borders = tcPr.find(qn('w:tcBorders'))
    if borders is None:
        borders = OxmlElement('w:tcBorders'); tcPr.append(borders)
    for edge in ('top','left','bottom','right'):
        if edge in kwargs:
            spec = kwargs[edge]
            el = borders.find(qn(f'w:{edge}'))
            if el is None:
                el = OxmlElement(f'w:{edge}'); borders.append(el)
            el.set(qn('w:val'), spec.get('val','single'))
            el.set(qn('w:sz'), str(spec.get('sz',4)))
            el.set(qn('w:color'), '000000')

# ======================================================================
# Word 渲染
# ======================================================================
def render_docx(spec, out_path, pack_name):
    from docx import Document
    from docx.oxml.ns import qn
    doc = Document()
    p = WORD_PACKS[pack_name]

    # 页边距
    for s in doc.sections:
        s.top_margin = Cm(2.54); s.bottom_margin = Cm(2.54)
        s.left_margin = Cm(3.17); s.right_margin = Cm(3.17)

    # 默认正文样式
    style = doc.styles['Normal']
    allow_color = p.get('allow_color', False)
    body_color = REF_BLACK_RGB if allow_color else GREY_BLACK
    style.font.name = p['body_font']
    style.font.size = Pt(p['body_size'])
    style.font.color.rgb = body_color
    style.element.rPr.rFonts.set(qn('w:eastAsia'), p['body_font'])
    pf = style.paragraph_format
    pf.line_spacing = p['line_spacing']
    pf.space_after = Pt(p['para_space_after'])

    def _paint(run, color):
        """着色：allow_color 用真实色，否则 _is_grey 守卫。"""
        if allow_color:
            run.font.color.rgb = color if color else REF_BLACK_RGB
        else:
            _force_run_black(run, color)

    # 标题
    meta = spec.get('meta', {})
    title = spec.get('title','')
    st = spec.get('subtitle','')
    h = doc.add_paragraph()
    hr = h.add_run(title)
    hr.bold = True; hr.font.size = Pt(p.get('title_size', 18)); hr.font.name = p['heading_font']
    hr._element.rPr.rFonts.set(qn('w:eastAsia'), p['heading_font'])
    _paint(hr, p.get('heading_color') if not allow_color else REF_BLACK_RGB)
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER
    if st:
        sp = doc.add_paragraph(); sp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sr = sp.add_run(st); sr.font.size = Pt(p.get('subtitle_size', 12)); sr.bold = allow_color
        _paint(sr, REF_TITLE_BLUE_RGB if allow_color else GREY_DARK)
    if meta.get('doc_no'):
        dp = doc.add_paragraph(); dp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        dr = dp.add_run(f"文号：{meta['doc_no']}"); _paint(dr, GREY_MID); dr.font.size=Pt(9)
    if meta.get('author') or meta.get('date'):
        ap = doc.add_paragraph(); ap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        ar = ap.add_run(f"{meta.get('author','')}　{meta.get('date','')}"); _paint(ar, body_color)

    # 表/图计数
    tbl_idx = 0; fig_idx = 0

    for sec in spec.get('sections', []):
        t = sec.get('type')
        if t == 'heading':
            lvl = sec.get('level',1)
            hp = doc.add_paragraph()
            r = hp.add_run(sec.get('text',''))
            r.bold = p.get('heading_bold', True)
            hs = {1: p.get('h1_size',16), 2: p.get('h2_size',14), 3:12, 4:11}.get(lvl, 12)
            r.font.size = Pt(hs)
            r.font.name = p['heading_font']
            r._element.rPr.rFonts.set(qn('w:eastAsia'), p['heading_font'])
            _paint(r, p.get('heading_color') if not allow_color else REF_BLACK_RGB)
            hp.paragraph_format.space_before = Pt(12); hp.paragraph_format.space_after = Pt(6)
        elif t == 'paragraph':
            tp = doc.add_paragraph()
            if p['indent_first_line']:
                tp.paragraph_format.first_line_indent = Cm(0.85)  # 约2字
            tr = tp.add_run(sec.get('text','')); _paint(tr, body_color)
        elif t == 'bullet':
            for it in sec.get('items',[]):
                bp = doc.add_paragraph(style='List Bullet')
                br = bp.add_run(it); _paint(br, body_color)
        elif t == 'numbered':
            for it in sec.get('items',[]):
                np_ = doc.add_paragraph(style='List Number')
                nr = np_.add_run(it); _paint(nr, body_color)
        elif t == 'table':
            tbl_idx += 1
            caption = sec.get('caption') or f"{p['caption_style'].replace('N',str(tbl_idx))}"
            cap_label = p['caption_style'].replace('N', str(tbl_idx))
            cap_p = doc.add_paragraph()
            cap_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            cr = cap_p.add_run(f"{cap_label}{p['caption_sep']}{caption.replace(cap_label,'').strip()}")
            _paint(cr, GREY_DARK); cr.font.size = Pt(10); cr.bold = True
            headers = sec.get('headers',[]); rows = sec.get('rows',[])
            kinds = sec.get('kinds', [])  # 数据性质矩阵（excel-ref/word-ref）
            ncol = max(len(headers), max((len(r) for r in rows), default=0))
            tbl = doc.add_table(rows=1+len(rows), cols=ncol)
            tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
            threeline = p.get('threeline_table', False)
            full_grid = p.get('table_full_grid', False)
            hdr_fill = p.get('header_fill', p.get('_header_fill','D9D9D9'))
            hdr_color = p.get('header_font_color', None) if allow_color else None
            hdr_size = p.get('header_size', 10)
            # 表头
            for j,htext in enumerate(headers):
                c = tbl.rows[0].cells[j]
                c.text = ''
                run = c.paragraphs[0].add_run(htext); run.bold=True
                run.font.size = Pt(hdr_size)
                if hdr_color and allow_color:
                    run.font.color.rgb = hdr_color
                else:
                    _paint(run, None)
                c.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
                if (hdr_fill or not threeline) and allow_color:
                    _force_cell_shade_raw(c, hdr_fill)
                elif not threeline:
                    _force_cell_shade(c, p.get('_header_fill','D9D9D9'))
            # 数据行
            for i,row in enumerate(rows):
                for j,val in enumerate(row):
                    c = tbl.rows[1+i].cells[j]
                    c.text=''
                    run = c.paragraphs[0].add_run(str(val))
                    run.font.size = Pt(hdr_size)
                    _paint(run, body_color)
                    # 数据性质底纹
                    if allow_color:
                        kind = kinds[i][j] if i < len(kinds) and j < len(kinds[i]) else None
                        fill = REF_KIND_FILL.get(kind) if kind else None
                        if fill:
                            _force_cell_shade_raw(c, fill)
            # 边框
            if threeline:
                # 三线表: 仅顶线(粗)+表头底线+底线(粗), 中间无线
                for j in range(ncol):
                    _set_cell_border(tbl.rows[0].cells[j], top={'sz':12}, bottom={'sz':6})
                    _set_cell_border(tbl.rows[-1].cells[j], bottom={'sz':12})
            else:
                # 全网格 0.5pt 黑（word-ref 与默认均走此分支）
                sz = 4 if full_grid else 4
                for r_ in tbl.rows:
                    for c_ in r_.cells:
                        _set_cell_border(c_, top={'sz':sz}, bottom={'sz':sz}, left={'sz':sz}, right={'sz':sz})
            # 题注脚注
            if sec.get('note'):
                np_ = doc.add_paragraph()
                nr = np_.add_run(f"{p['footnote_label']}{sec['note']}")
                _paint(nr, GREY_MID); nr.font.size = Pt(9)
        elif t == 'page_break':
            doc.add_page_break()

    doc.save(out_path)
    return out_path

# ======================================================================
# Excel 渲染
# ======================================================================
def _xl_border(style):
    side = Side(style=style if style!='threeline' else 'thin', color=XL_BLACK)
    return Border(left=side,right=side,top=side,bottom=side)

def render_xlsx(spec, out_path, pack_name):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = (spec.get('title','Sheet1'))[:31]
    p = EXCEL_PACKS[pack_name]

    # 单一表格: 取 sections 里第一个 table 渲染(多表则多 sheet)
    tables = [s for s in spec.get('sections',[]) if s.get('type')=='table']
    if not tables:
        # 无表格: 把 paragraphs 写成 A 列
        ws['A1'] = spec.get('title','')
        r = 2
        for s in spec.get('sections',[]):
            if s.get('type')=='paragraph':
                ws.cell(row=r, column=1, value=s.get('text',''))
                r += 1
        wb.save(out_path); return out_path

    first = True
    for ti, tspec in enumerate(tables):
        if not first:
            ws = wb.create_sheet(title=(tspec.get('caption') or f"表{ti+1}")[:31])
        first = False
        allow_color = p.get('allow_color', False)
        headers = tspec.get('headers',[]); rows = tspec.get('rows',[])
        kinds = tspec.get('kinds', [])        # 数据性质矩阵
        col_fmt = tspec.get('col_fmt', [])    # 每列 number_format
        col_align = tspec.get('col_align', [])  # 每列对齐: center/left/right
        col_widths = tspec.get('col_widths', [])
        ncol = max(len(headers), max((len(r) for r in rows), default=0))
        thin = Side(style='thin', color=XL_BLACK)
        full_border = Border(left=thin,right=thin,top=thin,bottom=thin)
        # ---- 标题行（行1）----
        if p.get('title_merge'):
            ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=max(ncol,1))
        ws.cell(row=1, column=1, value=tspec.get('caption', spec.get('title','')))
        c = ws.cell(row=1,column=1)
        c.font = Font(bold=True, size=p.get('title_size',14),
                      color=p.get('title_color', XL_BLACK))
        c.alignment = Alignment(horizontal='left', vertical='center')
        ws.row_dimensions[1].height = p.get('title_row_height', 22)
        # ---- 可选说明行（行2）----
        hrow = 2
        if p.get('desc_row') and tspec.get('desc'):
            ws.merge_cells(start_row=2, start_column=1, end_row=2, end_column=max(ncol,1))
            dc = ws.cell(row=2, column=1, value=tspec.get('desc'))
            dc.font = Font(size=11, color=XL_DARK)
            dc.alignment = Alignment(horizontal='left', vertical='top', wrap_text=True)
            ws.row_dimensions[2].height = 30
            hrow = 3
        # ---- 表头行 ----
        for j,htext in enumerate(headers):
            cell = ws.cell(row=hrow, column=j+1, value=htext)
            cell.font = Font(bold=p.get('header_bold',True),
                             color=p.get('header_font_color', XL_BLACK),
                             size=p.get('header_size',11),
                             name=p.get('data_font','Calibri') if allow_color else None)
            cell.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
            if p.get('header_fill'):
                cell.fill = PatternFill('solid', fgColor=p['header_fill'])
            if p.get('border') and p['border']!='threeline':
                cell.border = full_border if allow_color else _xl_border(p['border'])
        ws.row_dimensions[hrow].height = p.get('title_row_height', 22)
        # ---- 数据行 ----
        rh = p.get('row_height', None)
        for i,row in enumerate(rows):
            if rh: ws.row_dimensions[hrow+1+i].height = rh
            for j,val in enumerate(row):
                cell = ws.cell(row=hrow+1+i, column=j+1, value=_coerce(val))
                cell.font = Font(color=XL_BLACK, size=p.get('data_size',11),
                                 name=p.get('data_font') if allow_color else None)
                # 对齐
                al = col_align[j] if j < len(col_align) else None
                if al == 'center':
                    cell.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
                elif al == 'right' or (p.get('currency_right') and j>0 and _is_number(val)):
                    cell.alignment = Alignment(horizontal='right', vertical='center')
                else:
                    cell.alignment = Alignment(horizontal='left', vertical='top', wrap_text=True)
                # 边框
                if p.get('border') and p['border']!='threeline':
                    cell.border = full_border if allow_color else _xl_border(p['border'])
                # 数字格式
                if j < len(col_fmt) and col_fmt[j]:
                    cell.number_format = col_fmt[j]
                elif p.get('thousands_sep') and _is_number(val):
                    cell.alignment = Alignment(horizontal='right')
                    cell.number_format = '#,##0.00' if ('.' in str(val)) else '#,##0'
                # 数据性质底纹
                if allow_color:
                    kind = kinds[i][j] if i < len(kinds) and j < len(kinds[i]) else None
                    fill = p.get('kind_fill',{}).get(kind) if kind else None
                    if fill:
                        cell.fill = PatternFill('solid', fgColor=fill)
                # 条件色(非彩: 用粗体区分正负)
                if p.get('cond_color') and _is_number(val):
                    try:
                        fv = float(val)
                        if fv < 0: cell.font = Font(bold=True, color=XL_BLACK)
                    except: pass
        # 三线表: 去中间边框, 仅顶/表头底/底
        if p.get('border')=='threeline':
            top_side = Side(style='medium', color=XL_BLACK)
            mid_side = Side(style='thin', color=XL_BLACK)
            bot_side = Side(style='medium', color=XL_BLACK)
            for j in range(ncol):
                ws.cell(row=hrow, column=j+1).border = Border(top=top_side, bottom=mid_side)
                last = hrow+len(rows)
                ws.cell(row=last, column=j+1).border = Border(bottom=bot_side)
        # 列宽
        for j in range(ncol):
            if j < len(col_widths) and col_widths[j]:
                ws.column_dimensions[get_column_letter(j+1)].width = col_widths[j]
            else:
                maxlen = len(str(headers[j])) if j<len(headers) else 0
                for row in rows:
                    if j < len(row): maxlen = max(maxlen, len(str(row[j])))
                ws.column_dimensions[get_column_letter(j+1)].width = min(max(maxlen*1.2+2, 10), 40)
        # 冻结
        if p.get('freeze_first_row'):
            ws.freeze_panes = ws.cell(row=hrow+1, column=1)
        # 题注脚注（注释行）
        if tspec.get('note'):
            nr = hrow + len(rows) + 2
            nc = ws.cell(row=nr, column=1, value="备注")
            nc.font = Font(bold=True, size=11, color=XL_BLACK)
            nc2 = ws.cell(row=nr, column=2, value=tspec['note'])
            nc2.font = Font(size=11, color=XL_DARK)
            nc2.alignment = Alignment(horizontal='left', vertical='top', wrap_text=True)
        # 图例（仅 excel-ref，第一张表底部）
        if p.get('legend') and ti == 0:
            lr = hrow + len(rows) + 4
            legends = [
                ("黄", REF_INPUT_YELLOW, "输入：直接取自附件原值"),
                ("绿", REF_COMPUTE_GREEN, "公式：由公式/交叉计算得到"),
                ("橙", REF_CAUTION_ORANGE, "提示：口径/分歧/注意"),
            ]
            for k,(lab, hexv, desc) in enumerate(legends):
                lc = ws.cell(row=lr+k, column=1, value=lab)
                lc.fill = PatternFill('solid', fgColor=hexv)
                lc.font = Font(size=11, color=XL_BLACK)
                lc.alignment = Alignment(horizontal='center', vertical='center')
                lc.border = full_border
                dc = ws.cell(row=lr+k, column=2, value=desc)
                dc.font = Font(size=11, color=XL_DARK)

    wb.save(out_path)
    return out_path

def _is_number(v):
    try:
        float(str(v).replace(',','').replace('％','').replace('%',''))
        return True
    except: return False

def _coerce(v):
    """尝试把字符串数字转回数字类型, 便于 Excel 计算。"""
    s = str(v).strip() if v is not None else ''
    if s == '': return ''
    # 去千分位
    s2 = s.replace(',','')
    # 百分比
    m = re.match(r'^-?\d+(\.\d+)?%?$', s2)
    if m:
        try:
            if '%' in s:
                return round(float(s2.replace('%',''))/100, 4)
            if '.' in s2:
                return float(s2)
            return int(s2)
        except: return v
    return v

# ======================================================================
# Markdown 回退
# ======================================================================
def render_md(spec, out_path, pack_name):
    lines = []
    lines.append(f"# {spec.get('title','')}")
    if spec.get('subtitle'): lines.append(f"\n*{spec['subtitle']}*\n")
    tbl = 0
    for s in spec.get('sections',[]):
        t = s.get('type')
        if t=='heading': lines.append(f"\n{'#'*max(s.get('level',1),1)} {s.get('text','')}\n")
        elif t=='paragraph': lines.append(s.get('text','')+'\n')
        elif t=='bullet':
            for it in s.get('items',[]): lines.append(f"- {it}")
            lines.append('')
        elif t=='numbered':
            for i,it in enumerate(s.get('items',[]),1): lines.append(f"{i}. {it}")
            lines.append('')
        elif t=='table':
            tbl+=1
            cap = s.get('caption') or f"表{tbl}"
            lines.append(f"\n**{cap}**\n")
            headers=s.get('headers',[]); rows=s.get('rows',[])
            lines.append('| ' + ' | '.join(headers) + ' |')
            lines.append('|' + '|'.join(['---']*len(headers)) + '|')
            for r in rows:
                lines.append('| ' + ' | '.join(str(x) for x in r) + ' |')
            if s.get('note'): lines.append(f"\n*注：{s['note']}*\n")
        elif t=='page_break': lines.append('\n---\n')
    with open(out_path,'w',encoding='utf-8') as f:
        f.write('\n'.join(lines))
    return out_path

# ======================================================================
# CSV 渲染（复刻 关键数字核验表 风格：UTF-8 BOM, 七列, 原值不格式化）
# ======================================================================
def render_csv(spec, out_path, pack_name):
    import csv as _csv
    rows_out = []
    # 标题行（可选）
    if spec.get('title'):
        rows_out.append([spec.get('title','')])
    # 每个 table 输出为一组：表头 + 数据行
    for s in spec.get('sections', []):
        t = s.get('type')
        if t == 'table':
            headers = s.get('headers', [])
            rows = s.get('rows', [])
            # 原值写出：不经 _coerce, 不加千分位, 不四舍五入
            if headers:
                rows_out.append(headers)
            for r in rows:
                rows_out.append([('' if c is None else str(c)) for c in r])
            # 表间空行分隔
            rows_out.append([])
        elif t == 'paragraph' and s.get('text'):
            rows_out.append([s.get('text','')])
    # 写文件：utf-8-sig（带 BOM）, CRLF
    with open(out_path, 'w', encoding='utf-8-sig', newline='') as f:
        w = _csv.writer(f, quoting=_csv.QUOTE_MINIMAL, lineterminator='\r\n')
        for r in rows_out:
            w.writerow(r)
    return out_path

# ======================================================================
# PPT (精简: 标题+要点, 黑白)
# ======================================================================
def render_pptx(spec, out_path, pack_name):
    from pptx import Presentation
    from pptx.util import Pt, Inches
    from pptx.dml.color import RGBColor
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    # 封面
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
    tf = tb.text_frame; tf.word_wrap=True
    r = tf.paragraphs[0].add_run(); r.text = spec.get('title','')
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = RGBColor(0,0,0)
    if spec.get('subtitle'):
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = spec['subtitle']
        r2.font.size = Pt(20); r2.font.color.rgb = RGBColor(0x40,0x40,0x40)
    # 内容页: 每个 heading 起一页
    cur_slide = None; cur_tf = None
    for sec in spec.get('sections',[]):
        t = sec.get('type')
        if t=='heading':
            cur_slide = prs.slides.add_slide(blank)
            tb = cur_slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(1))
            r = tb.text_frame.paragraphs[0].add_run(); r.text = sec.get('text','')
            r.font.size = Pt(32); r.font.bold=True; r.font.color.rgb=RGBColor(0,0,0)
            bb = cur_slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5))
            cur_tf = bb.text_frame; cur_tf.word_wrap=True
        elif t=='paragraph' and cur_tf:
            p = cur_tf.add_paragraph(); r=p.add_run(); r.text=sec.get('text','')
            r.font.size=Pt(18); r.font.color.rgb=RGBColor(0,0,0)
        elif t=='bullet' and cur_tf:
            for it in sec.get('items',[]):
                p = cur_tf.add_paragraph(); p.level=1
                r=p.add_run(); r.text='• '+it; r.font.size=Pt(16); r.font.color.rgb=RGBColor(0x40,0x40,0x40)
    prs.save(out_path)
    return out_path

# ======================================================================
# 风格包轮换
# ======================================================================
def pick_pack(spec, out_dir, deliverable):
    if spec.get('style_pack') and spec['style_pack']!='auto':
        return spec['style_pack']
    # 轮换: 按 out/ 已有产物数
    import glob
    existing = glob.glob(os.path.join(out_dir, '*'))
    n = len(existing)
    if deliverable == 'xlsx':
        return ['excel-A','excel-B','excel-C'][n % 3]
    return ['word-A','word-B','word-C'][n % 3]

# ======================================================================
# 主
# ======================================================================
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('content_json')
    ap.add_argument('-o','--out', default='out')
    ap.add_argument('--pack', default=None)
    ap.add_argument('--format', default=None, help='覆盖 deliverable')
    ap.add_argument('--name', default=None, help='指定输出文件主名（不含扩展名），优先于 title；用于精确匹配 task_metadata.deliverable_files')
    args = ap.parse_args()

    with open(args.content_json, 'r', encoding='utf-8') as f:
        spec = json.load(f)

    os.makedirs(args.out, exist_ok=True)
    deliverable = (args.format or spec.get('deliverable','docx')).lower()
    if deliverable.startswith('.'): deliverable = deliverable[1:]

    pack = args.pack or pick_pack(spec, args.out, deliverable)

    base = args.name or spec.get('title','golden') or 'golden'
    base = re.sub(r'[\\/:*?"<>|]','_', base)[:80]
    if deliverable == 'docx':
        # 校验 pack 属于 word 族
        if pack not in WORD_PACKS: pack = pick_pack(spec, args.out,'docx')
        out_path = os.path.join(args.out, f"{base}.docx")
        render_docx(spec, out_path, pack)
    elif deliverable == 'xlsx':
        if pack not in EXCEL_PACKS: pack = pick_pack(spec, args.out,'xlsx')
        out_path = os.path.join(args.out, f"{base}.xlsx")
        render_xlsx(spec, out_path, pack)
    elif deliverable == 'pptx':
        out_path = os.path.join(args.out, f"{base}.pptx")
        render_pptx(spec, out_path, pack)
    elif deliverable == 'md':
        out_path = os.path.join(args.out, f"{base}.md")
        render_md(spec, out_path, pack)
    elif deliverable == 'csv':
        out_path = os.path.join(args.out, f"{base}.csv")
        render_csv(spec, out_path, pack)
    else:
        raise SystemExit(f"未知 deliverable: {deliverable}")

    print(json.dumps({
        "ok": True,
        "path": os.path.abspath(out_path),
        "pack": pack,
        "deliverable": deliverable,
    }, ensure_ascii=False))

if __name__ == '__main__':
    main()
